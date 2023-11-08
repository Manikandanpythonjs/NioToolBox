from tkinter import PhotoImage
import customtkinter
from PIL import Image
from CTkMessagebox import CTkMessagebox
from tkinter import filedialog, StringVar
import PyPDF2
from pptx import Presentation
from pptx.util import Inches


class PdftoPpt(customtkinter.CTk):
    def __init__(self):
        super().__init__()

        self.selectedpath = StringVar()

        def center_window(self, width, height):
            screen_width = self.winfo_screenwidth()
            screen_height = self.winfo_screenheight()
            x_coordinate = (screen_width - width) // 2
            y_coordinate = (screen_height - height) // 2
            self.geometry(f"{width}x{height}+{x_coordinate}+{y_coordinate}")

        self.title("pdftoppt")
        self._set_appearance_mode("light")
        center_window(self, 300, 200)
        self.resizable(False, False)
        self.configure(fg_color="white")
        self.attributes('-topmost', True)

        self.logotitle = customtkinter.CTkLabel(
            self, text="PDF TO PPTX", text_color="grey")
        self.logotitle.pack(padx=10, pady=10, anchor="n")

        self.openfile = customtkinter.CTkButton(
            self, text="Change PDF to PPTX ", command=self.open_file)
        self.openfile.pack(padx=20, pady=20)

        self.selectedpathLabel = customtkinter.CTkLabel(
            self,  text=" ", text_color="grey")
        self.selectedpathLabel.pack(padx=10, pady=10, anchor="n")

    def open_file(self):
        self.attributes('-topmost', False)
        self.file_path = filedialog.askopenfilename(
            filetypes=[("Pdf Files", "*.pdf")])
        if self.file_path:

            self.selectedpathLabel.configure(text=self.file_path)
            self.openfile.configure(
                text="Convert to PPTX", command=self.convertpdftoppt)
            self.selectedpath.set(self.file_path)
            print("Selected file:", self.file_path)
        self.attributes('-topmost', True)

    def convertpdftoppt(self):
        self.wm_attributes("-topmost", False)
        if (self.file_path):

            self.selectedpathLabel.configure(text=self.file_path)

            self.output_word_file = filedialog.asksaveasfilename(
                defaultextension=".docx", filetypes=[("Power Point Files", "*.pptx")])

            self.openfile.configure(
                text="Please Wait ....")
            prs = Presentation()

            with open(self.file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)

                for page_num in range(num_pages):
                    # Create a new slide for each page
                    slide = prs.slides.add_slide(prs.slide_layouts[5])

                    # Get the text from the page
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()

                    # Add the text to the slide
                    left = Inches(1)
                    top = Inches(1.5)
                    width = Inches(8)
                    height = Inches(5)
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = text

            prs.save(self.output_word_file)
            self.wm_attributes("-topmost", True)
            self.openfile.configure(
                text="Change PDF to PPTX", command=self.open_file)

            self.selectedpath.set("")

        else:

            pass
